

from pptx import Presentation


from pptx.shapes.graphfrm import GraphicFrame
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.chart.data import CategoryChartData


def remove_placeholder(slide, placeholder_text):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text == placeholder_text:
            shape.text_frame.clear()


def add_list_to_slide(prs, data_list):
    # Choose a slide layout (0 represents a blank slide)
    slide_layout = prs.slide_layouts[0]

    # Add a slide with the chosen layout
    # slide = prs.slides.add_slide(slide_layout)

    # Add a title to the slide
    title = slide.shapes.title

    # Add a text box for the list
    left, top, width, height = Inches(2), Inches(3.5), Inches(4), Inches(3)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame

    # Add each item in the list as a paragraph
    for item in data_list:
        p = text_frame.add_paragraph()
        p.text = f"• {item}"  # Use "•" for bullet points

    return prs
# read presentation from file
prs = Presentation('C:/Users/Lenovo/Desktop/Template/template.pptx')

# find the first chart object in the presentation
slideIdx = 0
for slide in prs.slides:

    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            print("Chart of type %s found in slide[%s, id=%s] shape[%s, id=%s, type=%s]"
                    % (chart.chart_type, slideIdx, slide.slide_id,
                       slide.shapes.index(shape), shape.shape_id, shape.shape_type ))
            break
    slideIdx += 1

# create list with changed category names
categorie_map = { 'Category 1': 'Saurabh', 'Category 2': 'List' ,'Category 3': 'List2' ,'Category 4': 'List4'}
new_categories = list(categorie_map[c] for c in chart.plots[0].categories)

print(new_categories)

# build new chart data with new category names and old data values
new_chart_data = CategoryChartData()
new_chart_data.categories = new_categories
print(chart.series)
data_list = ["Saurabh", "List", "List2", "List3"]
presentation = add_list_to_slide(prs, data_list)
for series in chart.series:


    new_chart_data.add_series(series.name,(34.5, 31.5,55,60))

# write the new chart data to the chart
chart.replace_data(new_chart_data)

# save everything in a new file
prs.save('chart.pptx')
