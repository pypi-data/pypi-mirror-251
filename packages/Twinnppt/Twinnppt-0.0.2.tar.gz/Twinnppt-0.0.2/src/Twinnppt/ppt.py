import json
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from pptx.chart.data import CategoryChartData

prs = Presentation('template.pptx')

def remove_picture_by_position(prs, slide_index, left, top):
    slide = prs.slides[slide_index]
    
    for shape in slide.shapes:
        if shape.left == left and shape.top == top and shape.shape_type == 13:  # Shape type 13 corresponds to Picture
            shape.element.getparent().remove(shape.element)

def add_list_to_slide(prs, data_list):
    # Choose a slide layout (0 represents a blank slide)
    slide_layout = prs.slide_layouts[0]
    for slide in prs.slides:

    # Add a slide with the chosen layout
    # slide = prs.slides.add_slide(slide_layout)

    # Add a title to the slide
        title = slide.shapes.title

        # Add a text box for the list
        left, top, width, height = Inches(2.5), Inches(1.5), Inches(4), Inches(3)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        # Add each item in the list as a paragraph
        for item in data_list:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"  # Use "•" for bullet points

    return prs

def create_chart(chart_type, data):
    if chart_type == "bar":
        plt.bar(data['labels'], data['values'])
    elif chart_type == "line":
        plt.plot(data['labels'], data['values'])
    elif chart_type == "pie":
        plt.pie(data['values'], labels=data['labels'], autopct='%1.1f%%')
    else:
        raise ValueError("Invalid chart type")

    plt.title(data.get('title', ''))
    plt.xlabel(data.get('xlabel', ''))
    plt.ylabel(data.get('ylabel', ''))

    
    # Save the chart to a BytesIO buffer
    buffer = BytesIO()
    plt.savefig(buffer, format='png')
    buffer.seek(0)
    plt.close()
    

    return buffer

def create_ppt(chart_image_buffer, ppt_file,data):
    prs = Presentation('template.pptx')
    slideIdx = 0
    for slide in prs.slides:

        for shape in slide.shapes:
            if shape.has_chart:
                chart = shape.chart
                print("Chart of type %s found in slide[%s, id=%s] shape[%s, id=%s, type=%s]"
                        % (chart.chart_type, slideIdx, slide.slide_id,
                           slide.shapes.index(shape), shape.shape_id, shape.shape_type ))
                break
        slideIdx += 1

    # Add a slide with a title and content layout
    # slide_layout = prs.slide_layouts[5]
    # slide = prs.slides.add_slide(slide_layout)

    # Add title and content to the slide
    # title = slide.shapes.title
    # title.text = "Chart Presentation"

    # Add the chart image to the slide
    right = Inches(6)
    top = Inches(1)

    
    # data_list = ["Saurabh", "List", "List2", "List3"]
    # presentation = add_list_to_slide(prs, data_list)
    remove_picture_by_position(prs, 0, right, top)
    pic = slide.shapes.add_picture(chart_image_buffer, right, top)
    # data_list = ["Saurabh", "List", "List2", "List3"]
    presentation = add_list_to_slide(prs, data['labels'])
    

    # Save the PowerPoint presentation
    prs.save(ppt_file)
    return ppt_file

def convert_ppt(input_json):
    data = json.loads(input_json)
    chart_type = data['chart_type']
    chart_data = data['data']

    chart_image_buffer = create_chart(chart_type, chart_data)

    ppt_file = "output.pptx"
    create_ppt(chart_image_buffer, ppt_file,chart_data)



# if __name__ == "__main__":
#     # Example JSON data
    

#     data = json.loads(input_json)
#     chart_type = data['chart_type']
#     chart_data = data['data']

#     chart_image_buffer = create_chart(chart_type, chart_data)

#     ppt_file = "output.pptx"
#     create_ppt(chart_image_buffer, ppt_file,chart_data)

#     print(f"Presentation saved to {ppt_file}")
